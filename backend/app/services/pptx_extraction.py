from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400
TITLE_TOP_LIMIT = int(1.7 * EMU_PER_INCH)
MAX_TITLE_CHARACTERS = 120


def extract_shape_text(shape: Any) -> str | None:
    if not getattr(shape, "has_text_frame", False):
        return None

    text = shape.text_frame.text.strip()
    return text or None


def normalize_text(value: str) -> str:
    return " ".join(value.split()).strip()


def normalize_for_compare(value: str) -> str:
    return normalize_text(value).lower()


def is_low_value_shape_text(value: str) -> bool:
    normalized = normalize_for_compare(value)

    if not normalized:
        return True

    if normalized.isdigit():
        return True

    return any(
        marker in normalized
        for marker in (
            "confidential",
            "oracle restricted",
            "all rights reserved",
            "copyright",
            "\u00a9",
        )
    )


def first_meaningful_line(value: str) -> str | None:
    for line in value.splitlines():
        normalized_line = normalize_text(line)
        if normalized_line:
            return normalized_line

    return None


def extract_placeholder_title(slide: Any) -> str | None:
    if slide.shapes.title is None:
        return None

    title = extract_shape_text(slide.shapes.title)
    return title if title and not is_low_value_shape_text(title) else None


def infer_title_from_text_shapes(text_shapes: list[dict[str, Any]]) -> str | None:
    candidates: list[tuple[int, int, str]] = []

    for text_shape in text_shapes:
        text = text_shape["text"]
        title = first_meaningful_line(text)

        if title is None or is_low_value_shape_text(title):
            continue

        if len(title) > MAX_TITLE_CHARACTERS:
            continue

        if text_shape["top"] > TITLE_TOP_LIMIT:
            continue

        candidates.append((text_shape["top"], text_shape["left"], title))

    if not candidates:
        return None

    candidates.sort(key=lambda candidate: (candidate[0], candidate[1]))
    return candidates[0][2]


def is_body_slide_text(value: str, slide_title: str | None) -> bool:
    if is_low_value_shape_text(value):
        return False

    if slide_title and normalize_for_compare(value) == normalize_for_compare(slide_title):
        return False

    return True


def extract_table_rows(shape: Any) -> list[list[str]] | None:
    if not getattr(shape, "has_table", False):
        return None

    rows: list[list[str]] = []

    for row in shape.table.rows:
        rows.append([cell.text.strip() for cell in row.cells])

    return rows


def extract_notes_text(slide: Any) -> str | None:
    if not getattr(slide, "has_notes_slide", False):
        return None

    notes_text = slide.notes_slide.notes_text_frame.text.strip()
    return notes_text or None


def save_shape_image(
    shape: Any,
    image_output_dir: Path,
    image_storage_prefix: str,
    slide_number: int,
    image_number: int,
) -> str | None:
    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
        return None

    image = shape.image
    image_extension = image.ext or "bin"
    image_filename = f"slide_{slide_number:03d}_image_{image_number:03d}.{image_extension}"
    image_path = image_output_dir / image_filename
    image_path.write_bytes(image.blob)
    return f"{image_storage_prefix}/{image_filename}"


def render_slide_text(slide: dict[str, Any]) -> str:
    lines = [f"Slide {slide['slide_number']}"]

    if slide["title"]:
        lines.append(f"Title: {slide['title']}")

    if slide["texts"]:
        lines.append("Text:")
        lines.extend(f"- {text}" for text in slide["texts"])

    if slide["tables"]:
        lines.append("Tables:")
        for table in slide["tables"]:
            lines.append(f"Table {table['index']}:")
            lines.extend(" | ".join(row) for row in table["rows"])

    if slide["notes"]:
        lines.append("Notes:")
        lines.append(slide["notes"])

    lines.append(f"Images: {slide['image_count']}")
    return "\n".join(lines)


def extract_pptx(
    file_path: Path,
    image_output_dir: Path,
    image_storage_prefix: str,
) -> dict[str, Any]:
    presentation = Presentation(file_path)
    image_output_dir.mkdir(parents=True, exist_ok=True)

    slides: list[dict[str, Any]] = []
    image_paths: list[str] = []
    table_count = 0
    total_image_count = 0

    for slide_index, slide in enumerate(presentation.slides, start=1):
        raw_text_shapes: list[dict[str, Any]] = []
        slide_tables: list[dict[str, Any]] = []
        slide_image_paths: list[str] = []
        slide_image_count = 0

        for shape in slide.shapes:
            shape_text = extract_shape_text(shape)
            if shape_text:
                raw_text_shapes.append(
                    {
                        "text": shape_text,
                        "left": getattr(shape, "left", 0),
                        "top": getattr(shape, "top", 0),
                    }
                )

            table_rows = extract_table_rows(shape)
            if table_rows is not None:
                table_count += 1
                slide_tables.append({"index": len(slide_tables) + 1, "rows": table_rows})

            image_path = save_shape_image(
                shape=shape,
                image_output_dir=image_output_dir,
                image_storage_prefix=image_storage_prefix,
                slide_number=slide_index,
                image_number=slide_image_count + 1,
            )
            if image_path:
                slide_image_count += 1
                total_image_count += 1
                image_paths.append(image_path)
                slide_image_paths.append(image_path)

        slide_title = extract_placeholder_title(slide) or infer_title_from_text_shapes(
            raw_text_shapes
        )
        slide_texts = [
            text_shape["text"]
            for text_shape in raw_text_shapes
            if is_body_slide_text(text_shape["text"], slide_title)
        ]

        slides.append(
            {
                "slide_number": slide_index,
                "title": slide_title,
                "texts": slide_texts,
                "tables": slide_tables,
                "notes": extract_notes_text(slide),
                "image_count": slide_image_count,
                "image_paths": slide_image_paths,
            }
        )

    return {
        "text_content": "\n\n".join(render_slide_text(slide) for slide in slides),
        "json_content": {
            "slide_count": len(slides),
            "slides": slides,
            "image_paths": image_paths,
            "table_count": table_count,
            "total_image_count": total_image_count,
        },
    }
